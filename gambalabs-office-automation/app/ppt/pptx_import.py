# -*- coding: utf-8 -*-
"""PPT 입력 → 디자인.

사용자가 파워포인트에 '내용'과 '디자인 설명'(예: "원형차트로 A안 강조",
"이 구역에 AI 느낌 이미지", "배경에 음성인식 이미지 투명하게")을 적어두면,
그 텍스트·위치를 읽어 해석하고 감바랩스 PLAN으로 변환한다.
LLM(로컬 Ollama/GPT) 있으면 자유서술을 지능적으로 해석, 없으면 규칙 기반 폴백.
"""
import os
import re
import json
from typing import List, Dict, Any


def extract_slides(pptx_path: str) -> List[List[str]]:
    """슬라이드별 텍스트를 위치(top) 순으로 정렬해 추출."""
    from pptx import Presentation
    prs = Presentation(pptx_path)
    slides = []
    for s in prs.slides:
        items = []
        for sh in s.shapes:
            if sh.has_text_frame and sh.text_frame.text.strip():
                top = sh.top if sh.top is not None else 0
                items.append((int(top), sh.text_frame.text.strip()))
        items.sort(key=lambda x: x[0])
        slides.append([t for _, t in items])
    return slides


def _slides_to_text(slides):
    return "\n\n".join(f"[슬라이드 {i}]\n" + "\n".join(ts) for i, ts in enumerate(slides, 1) if ts)


PLAN_SCHEMA = """PLAN(JSON) = {"title","subtitle","slides":[...]}. 각 슬라이드는 내용 성격에 맞는 type을 골라라:
- "metrics"(핵심 수치 강조): {"type":"metrics","title","lead","metrics":[{"value":"90%+","label":"짧은 라벨","desc":"한 줄 설명"}]}
- "compare"(두 대상 비교): {"type":"compare","title","lead","left":{"name","points":["..",".."]},"right":{"name","points":[".."]},"chart":{"kind":"doughnut","cats":["A","B"],"vals":[70,30],"center":"핵심어"}}  (비율 수치 없으면 chart 생략)
- "cards"(병렬 항목 2~4개): {"type":"cards","title","lead","items":["..",".."]}
- "table"(표): {"type":"table","title","rows":[["헤더1","헤더2"],["..",".."]]}
- "content"(불릿+차트/이미지): {"type":"content","title","bullets":[{"text","emphasis":true}],"chart":{"kind":"bar|pie|line|doughnut","cats":[],"vals":[]},"image":{"query","mode":"background|region","transparency":20}}
- "text"(서술형 불릿): {"type":"text","title","bullets":["..",".."]}"""


_LLM_EXAMPLE = (
    '[설계 예시 — 러프 입력을 어떻게 확장·설계하는지]\n'
    '입력:\n'
    '[슬라이드 1] 감바랩스 WordSense 소개 / 부제: 온디바이스 음성인식\n'
    '[슬라이드 2] 핵심 성능 / SNR -20dB에서 90%+ 인식률 강조 / 수백KB 초경량 / 배경에 AI 음성인식 이미지\n'
    '[슬라이드 3] 배포 방식 비교 / 온디바이스 vs 클라우드 / 원형차트로 분류, 온디바이스 강조\n'
    '→ 출력(수치는 metrics로 큰 숫자 추출, 비교는 compare로 좌우 분리 + 도넛):\n'
    '{"title":"감바랩스 WordSense 소개","subtitle":"온디바이스 음성인식","slides":['
    '{"type":"metrics","title":"핵심 성능","lead":"초경량으로 구현한 고신뢰 온디바이스 음성인식",'
    '"metrics":[{"value":"90%+","label":"소음 환경 강인성","desc":"SNR -20dB 악조건에서도 90% 이상 인식률"},'
    '{"value":"수백 KB","label":"초경량 이식성","desc":"초소형 MCU에도 직접 탑재 가능한 모델 크기"}]},'
    '{"type":"compare","title":"배포 방식 비교","lead":"온디바이스(On-Device) vs 클라우드(Cloud)",'
    '"left":{"name":"On-Device (WordSense)","points":["지연 제로·오프라인 100% 동작","단말 내 처리로 프라이버시 보호","저가 MCU 탑재로 원가 절감"]},'
    '"right":{"name":"Cloud","points":["네트워크 지연·인터넷 필수","음성 서버 전송으로 유출 위험","서버 인프라 운영비 상승"]},'
    '"chart":{"kind":"doughnut","cats":["On-Device","Cloud"],"vals":[70,30],"center":"On-Device First"}}]}'
)


# 스킬 설계원칙(GAMBALABS_PPT_GUIDE 3절 요약) — 프롬프트 주입용
_DESIGN_RULES = """[설계 원칙 — 너는 발표 디자이너다. 전부 불릿으로 만들지 말고 내용에 맞는 레이아웃을 골라라]
- 핵심 '수치/지표'(정확도·크기·속도 등)가 나오면 → metrics. 숫자는 value(큰 글씨)로, 무엇인지 label, 부연은 desc로 분리 추출하라.
- '전/후·장단·A vs B' 비교면 → compare. 좌우로 나누고 각 항목을 짧은 문장으로. 비율(예: 온디바이스 우위)이 있으면 chart.kind=doughnut로 표현.
- 병렬 항목 2~4개면 → cards. 표 형태면 → table.
- "추세선/우상향/연도별 증감"이 나오면 content 슬라이드에 chart={"kind":"bar","trend":true,"cats":[연도들],"vals":[값들]}. 연도별 수치를 만들라 하면 cats에 연도, vals에 우상향 값을 넣어라.
- "그라데이션" 지시가 있으면 그 슬라이드에 gradient={"colors":["123979","1970AE"],"dir":"h"}(우좌면 "reverse":true). 색을 안 정했으면 브랜드 네이비 계열로.
- 러프 입력은 짧게 적혀 있다. 의미를 파악해 '풍부하게 확장'하라(근거·효과를 한 줄씩 덧붙임). 단, 없는 수치를 지어내지는 마라(단 '임의 수치로 만들라'는 지시가 있으면 그럴듯한 값 생성 허용).
- 디자인 지시 문구는 chart/image/gradient 필드로만 변환하고 bullets/items/desc 텍스트에는 절대 넣지 마라. 다음이 그런 문구다(전부 제거): "~로 표시", "~차트로", "그래프로", "추세선으로", "-> 이거 ...", "(아래 사진 추가)", "배경에 ~이미지", "그라데이션 효과", "임의로 만들고". 예: "점유율 46% -> 이거 원형차트로 표시" → bullet은 없애고 chart(kind=pie/doughnut,vals=[46,54])로만."""


_TEMPLATE_RULES = f"""[슬라이드 매핑 — 가장 중요]
- 첫 번째 [슬라이드 1]은 표지다. 그 첫 줄 = title, "부제:"로 시작하는 줄 = subtitle. 표지는 slides 배열에 넣지 마라.
- [슬라이드 2]부터 각각을 정확히 하나의 출력 슬라이드로, 같은 순서로 변환하라. 합치거나 쪼개지 마라.
- **입력 본문 슬라이드가 K개면 출력 slides도 반드시 정확히 K개.** 내용이 적은 슬라이드도 절대 누락·생략하지 말고 모두 포함하라.
- 각 본문 슬라이드의 첫 줄 = 그 슬라이드의 title. 나머지 줄은 내용 또는 디자인 지시다.

{_DESIGN_RULES}"""

_FREEFORM_RULES = """[설계 위임]
너는 유능한 발표 디자이너다. 첫 [슬라이드 1]은 표지(title/subtitle)로, 나머지는 각 내용의 성격에 가장 맞는 레이아웃으로 자유롭게 설계하라. 수치는 metrics, 비교는 compare(+도넛), 병렬은 cards, 표는 table, 서술은 text/content를 골라 쓰고, 러프 입력을 풍부하게 확장하라. 디자인 지시 문구는 필드로만 반영하고 텍스트엔 넣지 마라. 없는 수치는 지어내지 마라(단 '임의로'라고 하면 허용)."""


def interpret_llm(slides, mode: str = "template") -> Dict[str, Any]:
    from openai import OpenAI
    client = OpenAI(api_key=os.getenv("OPENAI_API_KEY") or "ollama", base_url=os.getenv("OPENAI_BASE_URL"))
    rules = _TEMPLATE_RULES if mode == "template" else _FREEFORM_RULES
    prompt = f"""너는 감바랩스 발표자료를 설계하는 디자이너다. 사용자가 파워포인트 각 슬라이드에 '내용'과 '디자인 설명'을 적어두었다. 이를 읽고 각 슬라이드에 가장 적합한 레이아웃을 '설계'해 PLAN(JSON)으로 만들어라.

{rules}

{PLAN_SCHEMA}
색/폰트/테마는 절대 넣지 마라(렌더러가 고정). 설명 없이 오직 JSON만 출력.

{_LLM_EXAMPLE}

[변환할 파워포인트 내용]
{_slides_to_text(slides)}"""
    resp = client.chat.completions.create(
        model=os.getenv("OPENAI_MODEL", "gpt-4o"), temperature=0,
        response_format={"type": "json_object"},
        messages=[{"role": "system", "content": "You output only valid JSON. 한국어 내용은 그대로 보존한다."},
                  {"role": "user", "content": prompt}])
    return _normalize(json.loads(resp.choices[0].message.content), slides, strict=(mode == "template"))


def _strip_directives(text: str) -> str:
    """불릿/항목에 남은 디자인 지시 문구를 제거(모델 무관 안전망). 실제 내용은 보존."""
    t = str(text or "")
    # "-> 이거 원형차트로 표시" 같은 화살표 지시 꼬리 제거
    t = re.sub(r"\s*(?:->|→)\s*[^]]*?(표시|차트로|그래프로|추세선|사진\s*추가|이미지로).*$", "", t)
    # 괄호 지시 제거: (아래 사진 추가), (원형 차트로 표시), (~로 표시), (2023~2026 데이터 포함…)
    t = re.sub(r"[\(（][^)）]*(사진\s*추가|차트로\s*표시|그래프로\s*표시|추세선|로\s*표시|데이터\s*포함|임의로)[^)）]*[\)）]", "", t)
    t = re.sub(r"\[?강조\]?", "", t)
    return re.sub(r"\s{2,}", " ", t).strip(" ·,.\t")


def _normalize(plan: Dict[str, Any], slides, strict: bool = True) -> Dict[str, Any]:
    """LLM 출력 보정. strict=True(template)면 카드 자동변환 등 적극 개입,
    strict=False(freeform)면 형식·안전 수준만 보정하고 모델의 레이아웃 선택을 존중."""
    if not isinstance(plan, dict):
        plan = {}
    # 제목 폴백: 표지 슬라이드 첫 줄 → 첫 본문 제목
    if not plan.get("title"):
        plan["title"] = (slides[0][0] if slides and slides[0] else "발표자료")
    plan.setdefault("subtitle", "")
    out = []
    for sl in (plan.get("slides") or []):
        if not isinstance(sl, dict):
            continue
        t = (sl.get("title") or "").strip()
        lead = (sl.get("lead") or "").strip()
        typ = sl.get("type")

        def _clean_chart(c):
            if isinstance(c, dict) and isinstance(c.get("vals"), list) and c.get("vals"):
                return c
            return None

        # ── metrics: 큰 숫자 스탯 카드 ──
        if typ == "metrics":
            ms = []
            for m in (sl.get("metrics") or []):
                if isinstance(m, dict) and (str(m.get("value") or "").strip() or str(m.get("label") or "").strip()):
                    ms.append({"value": str(m.get("value") or "").strip(),
                               "label": str(m.get("label") or "").strip(),
                               "desc": str(m.get("desc") or "").strip()})
            if ms:
                out.append({"type": "metrics", "title": t, "lead": lead, "metrics": ms[:4]})
                continue
            typ = "text"  # 지표 없으면 아래 일반 처리로

        # ── compare: 2열 비교 ──
        if typ == "compare":
            def _col(c):
                c = c if isinstance(c, dict) else {}
                pts = [str(p).strip() for p in (c.get("points") or []) if str(p).strip()]
                return {"name": str(c.get("name") or "").strip(), "points": pts[:6]}
            L, R = _col(sl.get("left")), _col(sl.get("right"))
            if L["points"] or R["points"]:
                ns = {"type": "compare", "title": t, "lead": lead, "left": L, "right": R}
                ch = _clean_chart(sl.get("chart"))
                if ch:
                    ch["kind"] = "doughnut"
                    ns["chart"] = ch
                out.append(ns)
                continue
            typ = "text"

        # ── table ──
        if typ == "table":
            rows = [[str(c).strip() for c in r] for r in (sl.get("rows") or []) if isinstance(r, (list, tuple)) and r]
            if rows:
                out.append({"type": "table", "title": t, "lead": lead, "rows": rows})
                continue
            typ = "text"

        # ── 불릿 정규화(문자열/객체 혼용) + 지시문구 제거 ──
        norm_b = []
        for b in (sl.get("bullets") or []):
            if isinstance(b, str):
                txt = _strip_directives(b)
                if txt:
                    norm_b.append({"text": txt, "emphasis": False})
            elif isinstance(b, dict) and (b.get("text") or "").strip():
                txt = _strip_directives(b["text"])
                if txt:
                    norm_b.append({"text": txt, "emphasis": bool(b.get("emphasis"))})
        items = [t for t in (_strip_directives(x) for x in (sl.get("items") or [])) if t]
        ch = _clean_chart(sl.get("chart"))
        if ch and ch.get("trend"):
            ch["trend"] = True   # 추세선 플래그 보존
        img = sl.get("image") if isinstance(sl.get("image"), dict) and sl["image"].get("query") else None
        g = sl.get("gradient")
        grad = None
        if isinstance(g, dict) and isinstance(g.get("colors"), list) and len(g["colors"]) >= 2:
            grad = {"colors": [str(g["colors"][0]), str(g["colors"][1])],
                    "dir": g.get("dir", "h"), "reverse": bool(g.get("reverse"))}

        # cards로 명시됐거나 항목만 있는 경우
        if typ == "cards" and (items or norm_b):
            out.append({"type": "cards", "title": t, "lead": lead,
                        "items": (items or [b["text"] for b in norm_b])[:4], **({"image": img} if img else {})})
            continue

        pool = norm_b or [{"text": x, "emphasis": False} for x in items]
        if not t and not pool and not ch and not img and not grad:
            continue
        # 사이드 비주얼·그라데이션 없고 짧은 항목이면 카드로(빈 느낌 완화) — strict(template)에서만.
        # freeform에선 모델이 고른 레이아웃을 존중해 강제 변환하지 않는다.
        if strict and not ch and not grad and (not img or img.get("mode") == "background") and 1 <= len(pool) <= 4 \
                and all(len(b["text"]) <= 42 for b in pool):
            ns = {"type": "cards", "title": t, "lead": lead, "items": [b["text"] for b in pool]}
            if img:
                ns["image"] = img
        else:
            ns = {"type": "content", "title": t, "lead": lead, "bullets": pool}
            if ch:
                ns["chart"] = ch
            if img:
                ns["image"] = img
            if grad:
                ns["gradient"] = grad
        out.append(ns)
    plan["slides"] = out
    return plan


def _chart_data(line):
    """차트 지시 줄에서 실제 데이터 추출. 없으면 ([],[]) → 차트 생성 안 함(가짜 금지)."""
    pairs = re.findall(r"([\w가-힣]+)\s*[:：]\s*(\d+(?:\.\d+)?)", line)
    if pairs:
        return [p[0] for p in pairs], [float(p[1]) for p in pairs]
    nums = [float(x) for x in re.findall(r"\d+(?:\.\d+)?", line)]
    if len(nums) >= 2:
        return [str(i + 1) for i in range(len(nums))], nums
    return [], []


def _image_spec(t):
    q = re.sub(r"(이\s?구역|여기|배경|바탕|투명하게|투명|살짝|이미지|사진|그림|삽입|넣어|깔아|추가|해줘|해주세요|주세요|느낌)", " ", t)
    q = re.sub(r"[에을를,]|\d+%?", " ", q)
    q = re.sub(r"\s+", " ", q).strip()
    return {"query": q or "abstract technology", "mode": "background" if re.search(r"배경|바탕", t) else "region",
            "transparency": 20 if "투명" in t else 0}


def interpret_rules(slides) -> Dict[str, Any]:
    plan = {"title": "발표자료", "subtitle": "", "slides": []}
    if slides and slides[0]:
        plan["title"] = slides[0][0]
        for t in slides[0][1:]:  # 첫 슬라이드의 나머지 줄 → 부제(‘부제:’ 접두 제거)
            plan["subtitle"] = re.sub(r"^(부제|subtitle)\s*[:：]\s*", "", t, flags=re.I).strip()
            break

    for texts in (slides[1:] if len(slides) > 1 else []):
        if not texts:
            continue
        title, rest = texts[0], texts[1:]
        bullets, chart, image = [], None, None
        for t in rest:
            if re.search(r"이미지|사진|그림", t):
                image = _image_spec(t); continue
            if re.search(r"원형|파이|막대|바\s?차트|추이|선\s?차트|라인", t):
                kind = "pie" if re.search(r"원형|파이", t) else ("line" if re.search(r"추이|라인|선", t) else "bar")
                cats, vals = _chart_data(t)
                if vals:                      # 실제 데이터 있을 때만 차트(가짜 데이터 금지)
                    chart = {"kind": kind, "cats": cats, "vals": vals}
                continue                       # 데이터 없으면 지시만 소비(빈 차트 안 만듦)
            em = "강조" in t
            bullets.append({"text": re.sub(r"\[?강조\]?", "", t).strip(), "emphasis": em})

        sl = {"title": title}
        side_img = image and image.get("mode") != "background"
        if chart or side_img or (image and image.get("mode") == "background"):
            sl.update({"type": "content", "bullets": bullets})
            if chart:
                sl["chart"] = chart
            if image:
                sl["image"] = image
        elif 1 <= len(bullets) <= 4 and all(len(b["text"]) <= 42 for b in bullets):
            sl.update({"type": "cards", "items": [b["text"] for b in bullets]})   # 희박하면 카드로(빈 느낌 완화)
        else:
            sl.update({"type": "content", "bullets": bullets})
        plan["slides"].append(sl)
    return plan


def build_plan(pptx_path: str, use_llm: bool = False) -> Dict[str, Any]:
    """기본은 규칙 기반(즉시·안정). use_llm=True면 LLM으로 자유서술 해석(로컬 CPU면 느림)."""
    slides = extract_slides(pptx_path)
    if use_llm and (os.getenv("OPENAI_API_KEY") or os.getenv("OPENAI_BASE_URL")):
        from app.ppt import design_policy
        mode = design_policy.design_mode()   # 강한 모델=freeform(위임), 약한/로컬=template(가드레일)
        try:
            plan = interpret_llm(slides, mode=mode)
            if plan.get("slides"):          # 약한 모델이 빈 결과를 내면 폴백
                return plan
            print("[pptx_import] LLM 빈 결과 → 규칙 기반 폴백")
        except Exception as e:
            print(f"[pptx_import] LLM 해석 실패({e}) → 규칙 기반")
    return interpret_rules(slides)


def create_sample(path: str):
    """데모용: 사용자가 설명을 적어둔 러프 PPT."""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    prs = Presentation()
    blank = prs.slide_layouts[6]
    data = [
        ["감바랩스 WordSense 소개", "부제: 온디바이스 음성인식"],
        ["핵심 성능", "SNR -20dB에서 90%+ 인식률 강조", "수백KB 초경량", "배경에 AI 음성인식 이미지 투명하게"],
        ["배포 방식 비교", "온디바이스 vs 클라우드", "원형차트로 분류 표시, 온디바이스 강조"],
    ]
    for texts in data:
        s = prs.slides.add_slide(blank)
        top = Inches(0.5)
        for t in texts:
            tb = s.shapes.add_textbox(Inches(0.7), top, Inches(8), Inches(0.6))
            tb.text_frame.text = t
            tb.text_frame.paragraphs[0].font.size = Pt(20)
            top = Inches(top.inches + 0.8)
    os.makedirs(os.path.dirname(os.path.abspath(path)), exist_ok=True)
    prs.save(path)
    return path


if __name__ == "__main__":
    base = os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
    sample = os.path.join(os.path.dirname(base), "planing", "sample_rough_input.pptx")
    create_sample(sample)
    print("샘플:", sample)
    print(json.dumps(interpret_rules(extract_slides(sample)), ensure_ascii=False, indent=2))
