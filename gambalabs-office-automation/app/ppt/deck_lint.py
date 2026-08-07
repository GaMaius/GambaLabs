# -*- coding: utf-8 -*-
"""생성된 덱의 기하학적 결함을 코드로 잡아낸다.

모델이 쓴 pptxgenjs 코드는 '실행은 되는데 보기엔 엉망'인 경우가 많다.
사람이 스크린샷을 보내주기 전에 스스로 걸러내기 위한 검사기.
지적 사항은 문자열 목록으로 돌려주고, 그대로 모델에게 재수정 지시로 넘긴다.
"""
from typing import List

W_IN, H_IN = 13.333, 7.5
EMU = 914400.0

# freeform_helpers.js / freeform_engine.py 와 같은 그리드
HEADER_H = 1.05
MARGIN = 0.8
BODY_TOP = 1.93
BODY_BOTTOM = 6.88
FOOTER_TOP = 6.92          # 이 아래는 푸터(페이지 번호·구분선) — 본문 계산에서 뺀다


def _rects(slide):
    out = []
    for sh in slide.shapes:
        if sh.left is None or sh.top is None or sh.width is None or sh.height is None:
            continue
        kind = "pic" if (sh.shape_type == 13 or sh.__class__.__name__ == "Picture") else "shape"
        txt = ""
        try:
            if sh.has_text_frame:
                txt = sh.text_frame.text.strip()
                if txt:
                    kind = "text"
        except Exception:
            pass
        out.append({
            "kind": kind, "text": txt, "shape": sh,
            "x": sh.left / EMU, "y": sh.top / EMU,
            "w": sh.width / EMU, "h": sh.height / EMU,
        })
    return out


_NS_A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"


def _dark_outline(shape):
    """도형에 짙은 윤곽선이 그려져 있으면 그 색을 돌려준다.

    pptxgenjs는 line 옵션을 주지 않아도 <a:ln/>을 비워 두는데, 모델이
    line:{color:"000000"} 처럼 직접 테두리를 넣으면 카드가 촌스러워진다(실제 사례).
    """
    try:
        el = shape._element
    except AttributeError:
        return None
    for ln in el.iter(_NS_A + "ln"):
        for clr in ln.iter(_NS_A + "srgbClr"):
            v = (clr.get("val") or "").upper()
            if len(v) != 6:
                continue
            try:
                r, g, b = (int(v[i:i + 2], 16) for i in (0, 2, 4))
            except ValueError:
                continue
            if 0.299 * r + 0.587 * g + 0.114 * b < 130:
                return v
    return None


def _overlap(a, b):
    ox = min(a["x"] + a["w"], b["x"] + b["w"]) - max(a["x"], b["x"])
    oy = min(a["y"] + a["h"], b["y"] + b["h"]) - max(a["y"], b["y"])
    if ox <= 0 or oy <= 0:
        return 0.0
    return ox * oy


def _has_chart(slide) -> bool:
    for sh in slide.shapes:
        if getattr(sh, "has_chart", False):
            return True
    return False


def _has_picture(slide) -> bool:
    for sh in slide.shapes:
        if sh.shape_type == 13 or sh.__class__.__name__ == "Picture":
            return True
    return False


def lint(pptx_path: str, expect: dict = None) -> List[str]:
    """expect = {"n_slides":6, "charts":[4,5], "images":[3], "cover":True}

    기하 결함뿐 아니라 '기획안대로 만들어졌는지'까지 본다. 이게 없으면 모델이
    차트를 회색 네모+숫자로 그려 놓고도 검사를 통과한다(실제로 그랬다).
    """
    from pptx import Presentation
    prs = Presentation(pptx_path)
    problems = []
    expect = expect or {}
    # 실제 슬라이드 크기를 파일에서 읽는다(16:9 / 4:3 둘 다 지원).
    W_IN = round(prs.slide_width / EMU, 3)
    H_IN = round(prs.slide_height / EMU, 3)

    if len(prs.slides) == 0:
        return ["슬라이드가 하나도 없다."]

    n_exp = expect.get("n_slides")
    if n_exp and len(prs.slides) != n_exp:
        problems.append(f"슬라이드 수가 기획안과 다르다(기획 {n_exp}장, 생성 {len(prs.slides)}장).")

    for idx in expect.get("charts", []):
        if idx <= len(prs.slides) and not _has_chart(prs.slides[idx - 1]):
            problems.append(
                f"슬라이드 {idx}: 차트 슬라이드인데 실제 차트가 없다. "
                f"도형과 숫자로 흉내 내지 말고 addChart(s, 제목, chartData, THEME_INFO, prs, rect)를 호출하라.")

    for idx in expect.get("images", []):
        if idx <= len(prs.slides) and not _has_picture(prs.slides[idx - 1]):
            problems.append(f"슬라이드 {idx}: 배치해야 할 사진(IMG[{idx}])이 없다. addSlideImage로 넣어라.")

    for i, s in enumerate(prs.slides, 1):
        rs = _rects(s)
        if not rs:
            problems.append(f"슬라이드 {i}: 아무 요소도 없다.")
            continue

        # 1) 슬라이드 밖으로 나간 요소
        for r in rs:
            if r["x"] < -0.05 or r["y"] < -0.05 or \
               r["x"] + r["w"] > W_IN + 0.05 or r["y"] + r["h"] > H_IN + 0.05:
                label = (r["text"][:18] + "…") if r["text"] else r["kind"]
                problems.append(
                    f"슬라이드 {i}: '{label}'이(가) 슬라이드 밖으로 나감 "
                    f"(x={r['x']:.2f} y={r['y']:.2f} w={r['w']:.2f} h={r['h']:.2f}, 화면은 {W_IN}x{H_IN})")

        # 1-b) 표지가 아닌 본문 슬라이드의 세로 채움 — 이게 없으면 "위 1/3만 쓰고
        #      아래 절반이 백지"인 덱이 아무 지적 없이 통과한다(실제로 그랬다).
        if i > 1 or not expect.get("cover", True):
            core = [r for r in rs
                    if r["y"] >= 1.55 and r["y"] < FOOTER_TOP and r["h"] >= 0.2
                    and not (r["kind"] == "pic" and r["w"] > W_IN * 0.9)]   # 전면 배경 장식 제외
            if core:
                bottom = max(r["y"] + r["h"] for r in core)
                top = min(r["y"] for r in core)
                if bottom < BODY_BOTTOM - 0.7:
                    problems.append(
                        f"슬라이드 {i}: 본문이 y={bottom:.2f}에서 끝나 화면 아래 "
                        f"{(H_IN - bottom) / H_IN * 100:.0f}%가 비었다. "
                        f"카드/차트 높이를 키워 본문 박스(y {BODY_TOP}~{BODY_BOTTOM})를 채워라.")
                if top > BODY_TOP + 0.55:
                    problems.append(
                        f"슬라이드 {i}: 본문이 y={top:.2f}에서야 시작한다. 헤더 아래 여백이 너무 크다 "
                        f"(본문은 y={BODY_TOP}부터).")

        # 1-c) 카드에 짙은 테두리 — 파워포인트 기본 검은 윤곽선처럼 보여 촌스럽다
        for r in rs:
            if r["kind"] == "pic":
                continue
            col = _dark_outline(r["shape"])
            if col:
                problems.append(
                    f"슬라이드 {i}: 도형에 짙은 테두리(#{col})가 있다. "
                    f"모든 addShape에 line: {{ type: \"none\" }} 를 넣어 테두리를 없애라.")
                break

        # 2) 사진이 다른 요소에 가려짐
        #    장식은 제외한다: 헤더 밴드·로고(상단 1.1in 안), 전면 배경 이미지, 그라데이션 패널.
        #    이걸 안 빼면 밴드 위의 제목 텍스트를 '사진을 가림'으로 오탐해 루프가 헛돈다.
        pics = [r for r in rs
                if r["kind"] == "pic"
                and r["w"] * r["h"] > 0.5
                and r["y"] + r["h"] > 1.15                                  # 헤더 장식 제외
                and not (r["w"] > W_IN * 0.9 and r["h"] > (H_IN - 1.05) * 0.9)  # 전면 배경 제외
                and not (r["h"] > (H_IN - 1.05) * 0.9)]                     # 세로 전면 패널 제외
        for p_i, p in enumerate(pics):
            area = p["w"] * p["h"]
            covered = 0.0
            for r in rs[rs.index(p) + 1:]:      # 뒤에 그려진 것 = 위에 덮이는 것
                # 뒤에 그려진 '사진'도 앞 사진을 덮는다. 예전엔 사진끼리는 건너뛰어서
                # 전면 배경 사진이 그라데이션 패널을 통째로 가려도 통과했다.
                if r["kind"] == "pic" and r["w"] * r["h"] <= area:
                    continue
                covered += _overlap(p, r)
            if area and covered / area > 0.35:
                problems.append(
                    f"슬라이드 {i}: 사진이 카드/도형에 {covered / area * 100:.0f}% 가려졌다. "
                    f"사진 자리(x={p['x']:.2f} y={p['y']:.2f} w={p['w']:.2f} h={p['h']:.2f})를 비워라.")

        # 3) 본문 요소끼리 심하게 겹침
        body = [r for r in rs if r["y"] > 1.1 and r["kind"] != "pic"]
        for a_i in range(len(body)):
            for b_i in range(a_i + 1, len(body)):
                a, b = body[a_i], body[b_i]
                ov = _overlap(a, b)
                small = min(a["w"] * a["h"], b["w"] * b["h"])
                # 카드 위의 텍스트는 정상적인 포함 관계 → 완전히 감싸는 경우는 제외
                inside = (a["x"] <= b["x"] and a["y"] <= b["y"] and
                          a["x"] + a["w"] >= b["x"] + b["w"] and a["y"] + a["h"] >= b["y"] + b["h"]) or \
                         (b["x"] <= a["x"] and b["y"] <= a["y"] and
                          b["x"] + b["w"] >= a["x"] + a["w"] and b["y"] + b["h"] >= a["y"] + a["h"])
                if not inside and small and ov / small > 0.4:
                    la = (a["text"][:14] or a["kind"]); lb = (b["text"][:14] or b["kind"])
                    problems.append(f"슬라이드 {i}: '{la}'와 '{lb}'가 서로 겹친다.")

        # 4) 본문이 거의 비어 있음
        used = sum(r["w"] * r["h"] for r in rs if r["y"] > 1.1)
        if used < (W_IN * (H_IN - 1.05)) * 0.15:
            problems.append(f"슬라이드 {i}: 본문이 거의 비어 있다. 내용을 채우거나 배치를 키워라.")

        # 5) 상단에 제목 '글자'가 없다(표지 제외).
        #    밴드 이미지만 깔고 제목을 안 쓰는 경우가 있어 도형이 아니라 텍스트로 확인한다.
        if i > 1 or not expect.get("cover", True):
            if not any(r["kind"] == "text" and r["y"] < 1.2 for r in rs):
                problems.append(
                    f"슬라이드 {i}: 상단에 제목 텍스트가 없다. "
                    f"addHeader(s, \"제목\", THEME_INFO, ASSETS_DIR)로 제목을 넣어라.")

        # 6) 본문이 한쪽으로 쏠려 절반이 빈다(사진·그라데이션이 없는데도). 표지는 제외.
        has_side = any(r["kind"] == "pic" and r["y"] > 1.15 for r in rs)
        body_items = [] if (i == 1 and expect.get("cover", True)) else \
            [r for r in rs if r["y"] > 1.15]
        if body_items and not has_side:
            right = max(r["x"] + r["w"] for r in body_items)
            left = min(r["x"] for r in body_items)
            if right < W_IN * 0.62 or left > W_IN * 0.38:
                problems.append(
                    f"슬라이드 {i}: 본문이 한쪽에만 몰려 화면 절반이 비었다"
                    f"(가로 {left:.1f}~{right:.1f}in만 사용). 가용 폭 {MARGIN}~{W_IN - MARGIN:.2f}을 고루 쓰거나 "
                    f"가운데 정렬로 배치하라.")

        # 6-b) 좌우 여백이 안 맞음 — 왼쪽 0.8, 오른쪽 1.6처럼 되면 눈에 띄게 삐뚤어 보인다
        #      표지는 제외한다. 표지는 좌측 텍스트 + 우측 장식이 정상 구도다(오탐으로 루프가 헛돌았다).
        wide = [] if (i == 1 and expect.get("cover", True)) else [
                r for r in rs
                if r["y"] >= 1.55 and r["y"] < FOOTER_TOP and r["h"] >= 0.2 and r["w"] >= 0.5
                and not (r["kind"] == "pic" and r["w"] > W_IN * 0.9)]
        if wide:
            lm = min(r["x"] for r in wide)
            rm = W_IN - max(r["x"] + r["w"] for r in wide)
            if abs(lm - rm) > 0.45 and max(lm, rm) > 0.9:
                problems.append(
                    f"슬라이드 {i}: 좌우 여백이 다르다(왼쪽 {lm:.2f}in, 오른쪽 {rm:.2f}in). "
                    f"본문 폭은 x={MARGIN}에서 시작해 {W_IN - MARGIN:.2f}에서 끝나게 맞춰라.")

        # 7) 표지가 비었다
        if i == 1 and expect.get("cover", True):
            if len([r for r in rs if r["kind"] == "text"]) == 0:
                problems.append("슬라이드 1(표지)에 제목이 없다. addCover(prs, 제목, 부제, 발표자, THEME_INFO, ASSETS_DIR)를 쓰라.")

    # 중복 제거(같은 지적이 반복되면 프롬프트만 길어진다)
    seen, out = set(), []
    for p in problems:
        if p not in seen:
            seen.add(p)
            out.append(p)
    return out[:20]
